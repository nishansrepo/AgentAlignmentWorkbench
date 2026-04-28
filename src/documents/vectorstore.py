"""
documents/vectorstore.py - Format-aware ingestion, FAISS+BM25 hybrid retrieval.
Supports: .pdf .docx .pptx .xlsx .csv .html .md .json .txt
Graceful degradation: FAISS+BM25 -> BM25-only -> substring matching.
"""
from __future__ import annotations
import hashlib, json, logging, re
from datetime import datetime
from pathlib import Path
logger = logging.getLogger(__name__)
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
CHUNK_SIZE = 800
CHUNK_OVERLAP = 100

# === EXTRACTORS ===
def _extract_pdf(fp):
    try:
        import fitz
        doc = fitz.open(str(fp)); pages = []
        for i, p in enumerate(doc):
            t = p.get_text()
            if t.strip(): pages.append({"text": t, "metadata": {"page": i+1, "format": "pdf"}})
        doc.close(); return pages
    except: return []

def _extract_docx(fp):
    try:
        from docx import Document as D
        doc = D(str(fp)); blocks = []; hdr = "Document Start"
        for p in doc.paragraphs:
            t = p.text.strip()
            if not t: continue
            if p.style and p.style.name.startswith("Heading"): hdr = t; continue
            blocks.append({"text": t, "metadata": {"section": hdr, "format": "docx"}})
        return blocks
    except: return _extract_plaintext(fp)

def _extract_pptx(fp):
    try:
        from pptx import Presentation
        prs = Presentation(str(fp)); slides = []
        for i, sl in enumerate(prs.slides):
            texts = []
            for sh in sl.shapes:
                if sh.has_text_frame:
                    for p in sh.text_frame.paragraphs:
                        if p.text.strip(): texts.append(p.text.strip())
            if texts: slides.append({"text": "\n".join(texts), "metadata": {"slide": i+1, "format": "pptx"}})
        return slides
    except: return []

def _extract_spreadsheet(fp):
    try:
        import pandas as pd
        suf = fp.suffix.lower()
        sheets = {"Sheet1": pd.read_csv(fp, on_bad_lines="skip")} if suf == ".csv" else pd.read_excel(fp, sheet_name=None, engine="openpyxl")
        blocks = []
        for sn, df in sheets.items():
            blocks.append({"text": "Columns: " + " | ".join(str(c) for c in df.columns), "metadata": {"sheet": sn, "format": "spreadsheet"}})
            rows = []
            for idx, row in df.iterrows():
                rows.append(" | ".join(f"{c}: {v}" for c, v in row.items() if pd.notna(v)))
                if len(rows) >= 10:
                    blocks.append({"text": "\n".join(rows), "metadata": {"sheet": sn, "format": "spreadsheet"}}); rows = []
            if rows: blocks.append({"text": "\n".join(rows), "metadata": {"sheet": sn, "format": "spreadsheet"}})
        return blocks
    except: return _extract_plaintext(fp)

def _extract_html(fp):
    try:
        with open(fp, "r", encoding="utf-8", errors="replace") as f: raw = f.read()
        raw = re.sub(r"<script[^>]*>.*?</script>", "", raw, flags=re.DOTALL|re.IGNORECASE)
        raw = re.sub(r"<style[^>]*>.*?</style>", "", raw, flags=re.DOTALL|re.IGNORECASE)
        hp = re.compile(r"<(h[1-6])[^>]*>(.*?)</\1>", re.IGNORECASE|re.DOTALL)
        hs = list(hp.finditer(raw)); secs = []
        if hs:
            for i, m in enumerate(hs):
                s, e = m.end(), hs[i+1].start() if i+1 < len(hs) else len(raw)
                ht = re.sub(r"<[^>]+>", "", m.group(2)).strip()
                body = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", raw[s:e])).strip()
                if body: secs.append({"text": body, "metadata": {"section": ht, "format": "html"}})
        else:
            plain = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", raw)).strip()
            if plain: secs.append({"text": plain, "metadata": {"format": "html"}})
        return secs
    except: return _extract_plaintext(fp)

def _extract_markdown(fp):
    try:
        with open(fp, "r", encoding="utf-8", errors="replace") as f: text = f.read()
        hr = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)
        hs = list(hr.finditer(text))
        if not hs: return [{"text": text, "metadata": {"format": "markdown"}}]
        secs = []
        pre = text[:hs[0].start()].strip()
        if pre: secs.append({"text": pre, "metadata": {"section": "Preamble", "format": "markdown"}})
        for i, m in enumerate(hs):
            s, e = m.end(), hs[i+1].start() if i+1 < len(hs) else len(text)
            body = text[s:e].strip()
            if body: secs.append({"text": body, "metadata": {"section": m.group(2).strip(), "format": "markdown"}})
        return secs
    except: return _extract_plaintext(fp)

def _extract_json_file(fp):
    try:
        with open(fp, "r", encoding="utf-8") as f: data = json.load(f)
        blocks = []
        def _flat(obj, pfx=""):
            if isinstance(obj, dict):
                for k, v in obj.items(): _flat(v, f"{pfx}.{k}" if pfx else k)
            elif isinstance(obj, list):
                for i, item in enumerate(obj): _flat(item, f"{pfx}[{i}]")
            else:
                v = str(obj).strip()
                if v and len(v) > 10: blocks.append({"text": f"{pfx}: {v}", "metadata": {"key_path": pfx, "format": "json"}})
        _flat(data); return blocks
    except: return _extract_plaintext(fp)

_HEADER_PATTERNS = [
    re.compile(r"^(?:[0-9]+(?:\.[0-9]+)*\.?\s+|[IVXLC]+\.\s+)[A-Z][A-Za-z\s\-:&/,]+$", re.MULTILINE),
    re.compile(r"^[A-Z][A-Z\s\-:&/,]{4,60}$", re.MULTILINE),
    re.compile(r"^#{1,4}\s+.+$", re.MULTILINE),
    re.compile(r"^(?:Problem|Section|Part|Question|Chapter|Appendix)\s+[\dA-Za-z()]+[:.]\s*.*$", re.MULTILINE),
]

def _detect_headers(text):
    cands = set()
    for p in _HEADER_PATTERNS:
        for m in p.finditer(text):
            h = m.group().strip()
            if 3 <= len(h) <= 80: cands.add(h)
    ordered = [(text.find(h), h) for h in cands if text.find(h) >= 0]
    ordered.sort(); return [h for _, h in ordered]

def _extract_plaintext(fp):
    try:
        with open(fp, "r", encoding="utf-8", errors="replace") as f: text = f.read()
        if not text.strip(): return []
        headers = _detect_headers(text)
        if headers:
            blocks = []; last = 0; cur_h = "Preamble"
            for h in headers:
                idx = text.find(h, last)
                if idx >= 0:
                    c = text[last:idx].strip()
                    if c: blocks.append({"text": c, "metadata": {"section": cur_h, "format": "plaintext"}})
                    last = idx + len(h); cur_h = h
            rem = text[last:].strip()
            if rem: blocks.append({"text": rem, "metadata": {"section": cur_h, "format": "plaintext"}})
            if blocks: return blocks
        paras = [p.strip() for p in text.split("\n\n") if p.strip()]
        if not paras: return [{"text": text, "metadata": {"format": "plaintext"}}]
        blocks = []; cur = ""
        for p in paras:
            if len(cur) + len(p) + 2 <= CHUNK_SIZE: cur = cur + "\n\n" + p if cur else p
            else:
                if cur: blocks.append({"text": cur, "metadata": {"format": "plaintext"}})
                cur = p
        if cur: blocks.append({"text": cur, "metadata": {"format": "plaintext"}})
        return blocks
    except: return []

_EXTRACTORS = {".pdf": _extract_pdf, ".docx": _extract_docx, ".doc": _extract_docx, ".pptx": _extract_pptx,
    ".xlsx": _extract_spreadsheet, ".xls": _extract_spreadsheet, ".csv": _extract_spreadsheet,
    ".html": _extract_html, ".htm": _extract_html, ".md": _extract_markdown, ".json": _extract_json_file, ".txt": _extract_plaintext}

def extract_document(fp):
    return _EXTRACTORS.get(fp.suffix.lower(), _extract_plaintext)(fp)

# === CHUNKING ===
def _recursive_split(text, size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    if len(text) <= size: return [text] if text.strip() else []
    for sep in ["\n\n", "\n", ". ", " "]:
        parts = text.split(sep)
        if len(parts) > 1: return _merge_parts(parts, sep, size, overlap)
    return [text[i:i+size] for i in range(0, len(text), size - overlap)]

def _merge_parts(parts, sep, size, overlap):
    chunks = []; cur = ""
    for p in parts:
        p = p.strip()
        if not p: continue
        cand = cur + sep + p if cur else p
        if len(cand) <= size: cur = cand
        else:
            if cur: chunks.append(cur)
            if len(p) > size: chunks.extend(_recursive_split(p, size, overlap)); cur = ""
            else:
                if chunks and overlap > 0: cur = chunks[-1][-overlap:] + sep + p
                else: cur = p
    if cur and cur.strip(): chunks.append(cur)
    return chunks

def chunk_blocks(blocks, filename, rationale):
    chunks = []; idx = 0
    for block in blocks:
        text = block["text"]; meta = block.get("metadata", {})
        subs = [text] if len(text) <= CHUNK_SIZE else _recursive_split(text)
        for sub in subs:
            if not sub.strip(): continue
            cid = hashlib.md5(f"{filename}_{idx}_{sub[:40]}".encode()).hexdigest()[:10]
            chunks.append({"chunk_id": f"{filename}_{idx:03d}_{cid}", "text": sub,
                "metadata": {"filename": filename, "rationale": rationale[:300],
                    "section": meta.get("section", meta.get("title", "")),
                    "page": meta.get("page", ""), "format": meta.get("format", "unknown"),
                    "chunk_index": idx, "ingested_at": datetime.now().isoformat()}})
            idx += 1
    return chunks

# === STORE ===
class DocumentStore:
    def __init__(self):
        self._chunks = []; self._faiss_index = None; self._embeddings = None; self._bm25 = None; self._use_faiss = False
        try:
            from langchain_huggingface import HuggingFaceEmbeddings
            self._embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
            self._use_faiss = True; logger.info("FAISS mode: %s", EMBEDDING_MODEL)
        except: logger.info("Using BM25-only mode.")

    @property
    def count(self): return len(self._chunks)
    @property
    def supported_formats(self): return sorted(_EXTRACTORS.keys())

    def ingest_file(self, filepath, rationale=""):
        if not filepath.exists(): return 0
        blocks = extract_document(filepath)
        if not blocks: return 0
        chunks = chunk_blocks(blocks, filepath.name, rationale)
        if not chunks: return 0
        self._chunks.extend(chunks)
        if self._use_faiss:
            try: self._rebuild_faiss()
            except: self._use_faiss = False
        self._rebuild_bm25()
        return len(chunks)

    def _rebuild_faiss(self):
        from langchain_community.vectorstores import FAISS
        from langchain.schema import Document
        docs = [Document(page_content=c["text"], metadata=c["metadata"]) for c in self._chunks]
        self._faiss_index = FAISS.from_documents(docs, self._embeddings)

    def _rebuild_bm25(self):
        try:
            from rank_bm25 import BM25Okapi
            tok = [c["text"].lower().split() for c in self._chunks]
            if tok: self._bm25 = BM25Okapi(tok)
        except: self._bm25 = None

    def query(self, query_text, n_results=4):
        if not self._chunks: return []
        dense = self._qf(query_text, n_results*2); sparse = self._qb(query_text, n_results*2)
        if dense and sparse: return self._rrf(dense, sparse, n_results)
        elif dense: return dense[:n_results]
        elif sparse: return sparse[:n_results]
        else: return self._qs(query_text, n_results)

    def _qf(self, q, k):
        if not self._use_faiss or not self._faiss_index: return []
        try:
            docs = self._faiss_index.similarity_search_with_score(q, k=min(k, len(self._chunks)))
            return [{"chunk_id": d.metadata.get("filename","")+"_f", "text": d.page_content,
                "filename": d.metadata.get("filename",""), "section": d.metadata.get("section",""),
                "rationale": d.metadata.get("rationale",""), "score": round(float(s),4)} for d,s in docs]
        except: return []

    def _qb(self, q, k):
        if not self._bm25: return []
        try:
            scores = self._bm25.get_scores(q.lower().split())
            top = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:k]
            return [{"chunk_id": self._chunks[i]["chunk_id"], "text": self._chunks[i]["text"],
                "filename": self._chunks[i]["metadata"]["filename"], "section": self._chunks[i]["metadata"].get("section",""),
                "rationale": self._chunks[i]["metadata"].get("rationale",""), "score": round(float(scores[i]),4)}
                for i in top if scores[i] > 0]
        except: return []

    def _qs(self, q, k):
        words = [w for w in q.lower().split() if len(w) > 2]
        scored = [(sum(c["text"].lower().count(w) for w in words), c) for c in self._chunks]
        scored.sort(key=lambda x: x[0], reverse=True)
        return [{"text": c["text"], "filename": c["metadata"]["filename"], "section": c["metadata"].get("section",""),
            "rationale": c["metadata"].get("rationale",""), "score": s} for s, c in scored[:k] if s > 0]

    def _rrf(self, dense, sparse, n, k=60):
        scores = {}; items = {}
        for r, item in enumerate(dense):
            cid = item.get("chunk_id", item["text"][:30])
            scores[cid] = scores.get(cid, 0) + 1/(k+r+1); items[cid] = item
        for r, item in enumerate(sparse):
            cid = item.get("chunk_id", item["text"][:30])
            scores[cid] = scores.get(cid, 0) + 1/(k+r+1)
            if cid not in items: items[cid] = item
        ranked = sorted(scores.items(), key=lambda x: x[1], reverse=True)
        return [{**items[cid], "rrf_score": round(rrf, 4)} for cid, rrf in ranked[:n]]

    def get_all_context(self, max_chunks=8):
        parts = []
        for c in self._chunks[:max_chunks]:
            m = c["metadata"]
            h = f"[From: {m.get('filename','?')} | Section: {m.get('section','?')} | Purpose: {m.get('rationale','N/A')}]"
            parts.append(f"{h}\n{c['text']}")
        return "\n\n---\n\n".join(parts)

    def get_ingestion_summary(self):
        files = {}
        for c in self._chunks:
            fn = c["metadata"]["filename"]
            if fn not in files: files[fn] = {"filename": fn, "format": c["metadata"].get("format","?"), "rationale": c["metadata"].get("rationale",""), "chunk_count": 0, "sections": set()}
            files[fn]["chunk_count"] += 1
            s = c["metadata"].get("section","")
            if s: files[fn]["sections"].add(s)
        return [{"filename": f["filename"], "format": f["format"], "rationale": f["rationale"][:80],
            "chunks": f["chunk_count"], "sections": len(f["sections"])} for f in files.values()]

    def clear(self):
        self._chunks = []; self._faiss_index = None; self._bm25 = None
